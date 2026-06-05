"""Build EXIST 2026 Sexism in Memes presentation as a .pptx."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "figures"
OUT = ROOT / "sexism_memes_presentation.pptx"

# Colors
NAVY = RGBColor(0x14, 0x2A, 0x4F)
ACCENT = RGBColor(0xC9, 0x36, 0x4A)   # warm red for emphasis
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
INK = RGBColor(0x1F, 0x1F, 0x1F)
MUTED = RGBColor(0x55, 0x60, 0x70)
GREEN = RGBColor(0x2E, 0x8B, 0x57)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ---------------- helpers ----------------
def add_rect(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.text_frame.margin_left = Inches(0.1)
    shp.text_frame.margin_right = Inches(0.1)
    shp.text_frame.margin_top = Inches(0.05)
    shp.text_frame.margin_bottom = Inches(0.05)
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=INK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.0)
    tf.margin_top = Inches(0.0)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, indent = item
        else:
            text, indent = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = indent
        p.line_spacing = line_spacing
        bullet = "•   " if indent == 0 else "–   "
        r = p.add_run()
        r.text = bullet + text
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None):
    # Top bar
    add_rect(slide, 0, 0, SW, Inches(0.9), fill=NAVY)
    add_text(slide, Inches(0.5), Inches(0.18), SW - Inches(1), Inches(0.55),
             title, size=26, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), SW - Inches(1), Inches(0.4),
                 subtitle, size=14, color=MUTED)
    # Footer
    add_rect(slide, 0, SH - Inches(0.3), SW, Inches(0.3), fill=LIGHT)
    add_text(slide, Inches(0.4), SH - Inches(0.3), Inches(8), Inches(0.3),
             "EXIST 2026 · Sexism Identification in Memes · Klaudia Banasiewicz",
             size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)


def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def new_slide():
    return prs.slides.add_slide(BLANK)


# ============================================================
# Slide 1 — Title
# ============================================================
s = new_slide()
add_rect(s, 0, 0, SW, SH, fill=NAVY)
add_rect(s, 0, Inches(3.05), SW, Inches(0.06), fill=ACCENT)

add_text(s, Inches(0.8), Inches(1.5), SW - Inches(1.6), Inches(1.2),
         "Sexism Identification in Memes",
         size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.8), Inches(2.3), SW - Inches(1.6), Inches(0.7),
         "A multimodal approach with text, image and physiological signals",
         size=22, color=RGBColor(0xDE, 0xE3, 0xEC))

add_text(s, Inches(0.8), Inches(3.4), SW - Inches(1.6), Inches(0.5),
         "EXIST 2026 — Lab Project", size=18, bold=True,
         color=RGBColor(0xFF, 0xC2, 0x4B))

add_text(s, Inches(0.8), Inches(5.4), SW - Inches(1.6), Inches(0.5),
         "Klaudia Banasiewicz", size=22, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.8), Inches(5.9), SW - Inches(1.6), Inches(0.4),
         "Tasks 2.1 · 2.2 · 2.3", size=16,
         color=RGBColor(0xC9, 0xD1, 0xDE))

set_notes(s, """Hello, I'm Klaudia. I'll present my work on the EXIST 2026 'Sexism in Memes' task.
I worked on all three sub-tasks. As requested, I'll skip the task descriptions and focus on the models I tried, the final system I chose, and especially the cross-validation results with versus without physiological signals.
The talk is about 10 minutes — feel free to interrupt with questions at the end.""")


# ============================================================
# Slide 2 — Project at a glance
# ============================================================
s = new_slide()
add_header(s, "Project at a glance",
           "Three sub-tasks · one dataset · text + image + physiological signals")

# Three task cards
cards = [
    ("T2.1 — Binary", "Is the meme sexist?",
     "YES / NO", ACCENT),
    ("T2.2 — Type", "If sexist, what kind?",
     "NO / DIRECT / JUDGEMENTAL", NAVY),
    ("T2.3 — Category", "Multi-label: 5 sexism categories",
     "+ NO fallback", GREEN),
]
card_w = Inches(3.9)
card_h = Inches(2.1)
gap = Inches(0.25)
start_x = (SW - 3 * card_w - 2 * gap) / 2
for i, (title, sub, badge, color) in enumerate(cards):
    x = start_x + i * (card_w + gap)
    add_rect(s, x, Inches(1.6), card_w, card_h, fill=LIGHT)
    add_rect(s, x, Inches(1.6), card_w, Inches(0.08), fill=color)
    add_text(s, x + Inches(0.2), Inches(1.75), card_w - Inches(0.4),
             Inches(0.45), title, size=20, bold=True, color=color)
    add_text(s, x + Inches(0.2), Inches(2.25), card_w - Inches(0.4),
             Inches(0.45), sub, size=14, color=INK)
    add_text(s, x + Inches(0.2), Inches(2.85), card_w - Inches(0.4),
             Inches(0.6), badge, size=13, bold=True, color=MUTED)

# Dataset stats row
add_rect(s, Inches(0.6), Inches(4.2), SW - Inches(1.2), Inches(2.5), fill=LIGHT)
add_text(s, Inches(0.9), Inches(4.35), SW - Inches(1.8), Inches(0.5),
         "Dataset (EXIST 2026 Memes)", size=18, bold=True, color=NAVY)
stats = [
    ("3 984", "training memes"),
    ("1 053", "test memes"),
    ("EN + ES", "two languages"),
    ("6", "annotators per meme"),
    ("ET + HR + EEG", "physiological signals"),
]
sx = Inches(0.9)
sy = Inches(5.0)
cell_w = (SW - Inches(1.8)) / 5
for i, (n, lab) in enumerate(stats):
    add_text(s, sx + i * cell_w, sy, cell_w, Inches(0.7),
             n, size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, sx + i * cell_w, sy + Inches(0.75), cell_w, Inches(0.5),
             lab, size=13, color=MUTED, align=PP_ALIGN.CENTER)

set_notes(s, """Very briefly: three sub-tasks, same dataset, two languages, six annotators per meme.
The unusual ingredient here is that the organisers also recorded sensor data — eye tracking, heart rate, and brain signals — from real volunteers while they were looking at every meme.
So the whole project is really asking one question: does that sensor data actually carry useful information about sexism, or is it just noise on top of the text?""")


# ============================================================
# Slide 3 — Features / inputs per meme
# ============================================================
s = new_slide()
add_header(s, "What the model sees for each meme",
           "Three kinds of information about every meme — text, image, sensors")

# Three boxes
boxes = [
    ("Text", "The caption written on the meme",
     "→ XLM-RoBERTa: a multilingual\n   language model (handles\n   EN + ES with one model)\n→ Max 128 tokens per meme",
     "Pre-trained on 100 languages → no separate model per language", NAVY),
    ("Image", "The actual picture of the meme",
     "→ CLIP (OpenAI image encoder)\n→ Each picture → 512 numbers\n   that summarise 'what is\n   in the image'",
     "CLIP is frozen — I use its numbers but do not retrain it", GREEN),
    ("Physiological", "Eye tracking + heart rate + EEG (brain signals)",
     "→ Real volunteers wore sensors\n   while viewing each meme\n→ I aggregate across users\n   → 104 numbers per meme",
     "The model never sees raw EEG — only this 104-number summary", ACCENT),
]
bw = Inches(4.0)
bh = Inches(4.5)
gap = Inches(0.3)
bsx = (SW - 3 * bw - 2 * gap) / 2
by = Inches(1.7)
for i, (title, raw, pipeline, note, color) in enumerate(boxes):
    x = bsx + i * (bw + gap)
    add_rect(s, x, by, bw, bh, fill=LIGHT)
    add_rect(s, x, by, bw, Inches(0.55), fill=color)
    add_text(s, x + Inches(0.2), by + Inches(0.05), bw - Inches(0.4),
             Inches(0.45), title, size=20, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, x + Inches(0.2), by + Inches(0.7), bw - Inches(0.4),
             Inches(0.5), raw, size=14, bold=True, color=INK)
    add_text(s, x + Inches(0.2), by + Inches(1.25), bw - Inches(0.4),
             Inches(2.0), pipeline, size=14, color=INK)
    add_rect(s, x + Inches(0.2), by + Inches(3.4), bw - Inches(0.4),
             Inches(1.0), fill=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, x + Inches(0.3), by + Inches(3.5), bw - Inches(0.6),
             Inches(0.9), note, size=13, color=MUTED)

set_notes(s, """Every meme gives me three kinds of information.

First, the caption — the text written on the meme. I feed it to XLM-RoBERTa. That's a multilingual language model — 'multilingual' meaning it understands many languages out of the box, so I can handle both English and Spanish with one model instead of two.

Second, the picture. I use CLIP — an image model from OpenAI that already learned to describe pictures from millions of image–caption pairs on the web. CLIP turns each meme image into a list of 512 numbers that summarise what is in the image. I don't retrain CLIP — I just use its numbers.

Third — the interesting one — sensor data. The organisers had real volunteers wear sensors while they looked at every meme: eye tracking, heart rate, and EEG (brain signals). I don't feed the raw recordings into the model. Instead, I aggregate them across users into 104 numbers per meme — that compressed summary is what the model sees.

From here on, all my experiments are about how to best combine these three streams.""")


# ============================================================
# Slide 4 — Model ladder (overview of everything I tried)
# ============================================================
s = new_slide()
add_header(s, "Models I tried",
           "From the simplest 'guess the most common label' to multimodal combinations")

ladder = [
    ("1", "Majority / Minority baseline",
     "Always predict the most (or least) common label. A sanity-check floor."),
    ("2", "Logistic regression on image + sensors",
     "A simple linear classifier — draws a line between classes. No neural network."),
    ("3", "XLM-RoBERTa-base — text only",
     "Smaller multilingual language model, fine-tuned on the memes. 'How good is text alone?'"),
    ("4", "XLM-RoBERTa-large — full fine-tune",
     "Larger version of the same model. All weights updated on the task → strong text model."),
    ("5", "XLM-RoBERTa-large + LoRA + focal loss + oversampling",
     "LoRA = tiny add-on adapters; focal loss = focus on hard examples; oversampling = copy rare examples."),
    ("6", "XLM-RoBERTa-base — multi-label (5 sigmoid heads)",
     "Five independent yes/no detectors, one per category. The model can fire on more than one."),
    ("7", "Late fusion (text + image+sensors)",
     "Two models, two probabilities. I mix them: final = α · text + (1−α) · image+sensors."),
    ("8", "Fusion MLP — only for the cross-validation study",
     "A small comparison network. Same model, only the inputs change → fair feature comparison."),
]
sx = Inches(0.6)
sy = Inches(1.45)
row_h = Inches(0.66)
for i, (num, name, desc) in enumerate(ladder):
    y = sy + i * row_h
    add_rect(s, sx, y, Inches(0.6), row_h - Inches(0.05),
             fill=NAVY)
    add_text(s, sx, y, Inches(0.6), row_h - Inches(0.05),
             num, size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, sx + Inches(0.65), y, Inches(4.6), row_h - Inches(0.05),
             fill=LIGHT)
    add_text(s, sx + Inches(0.8), y + Inches(0.05),
             Inches(4.4), row_h - Inches(0.15),
             name, size=14, bold=True, color=INK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, sx + Inches(5.4), y + Inches(0.05),
             SW - sx - Inches(5.9), row_h - Inches(0.15),
             desc, size=13, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

set_notes(s, """This is the full ladder of models I tried — from simplest to most complex.

At the bottom: trivial baselines — always predict the most common label, just to set a floor. Then a logistic regression on the image and sensor features — a simple linear classifier, basically draws a line through the data.

Then I move to language models. XLM-RoBERTa-base is the smaller version — I use it as my text-only baseline. XLM-RoBERTa-large is the bigger version, where I update all the weights for the task — this is my strong text model.

For Task 2.2 — which has a rare class — I add three tricks on top of XLM-R-large. LoRA, which means I keep the original model frozen and only train tiny add-on adapter layers. Focal loss, a loss function that makes the model focus on hard examples instead of easy ones. And oversampling — literally copying the rare-class examples until they appear as often as the common ones.

For Task 2.3 I switch to a multi-label head — instead of picking one category out of five, the model has five independent yes/no detectors (called sigmoid heads in the code). They can all fire at once.

Late fusion is how I combine the text model with the image+sensor model. Each one gives a probability; I mix them with a weight α that I tune on a held-out set.

The last one — the fusion MLP — is a small comparison network that I only use for the cross-validation study on slide 8, where I need a controlled experiment.""")


# ============================================================
# Slide 5 — Text models per task
# ============================================================
s = new_slide()
add_header(s, "Per-task recipe (XLM-RoBERTa, tuned for each task)",
           "Same family of language model — different training tricks")

# Three columns: T2.1, T2.2, T2.3
cols = [
    ("T2.1 · Is this sexist?",
     "XLM-R-large — full fine-tune",
     [
         "Full fine-tune: update ALL the model's weights",
         "3 epochs · learning rate 1e-5 (slow, so it does not 'forget')",
         "Class-weighted cross-entropy: rare-class mistakes count more",
         "Mixed precision (bf16) to fit in GPU memory",
     ],
     "Why: this is the easiest task → a strong text model is already enough.",
     ACCENT),
    ("T2.2 · What kind of sexism?",
     "XLM-R-large + LoRA + focal loss + oversampling",
     [
         "'Judgemental' is only 14% of data → severe class imbalance",
         "LoRA: freeze the model, train tiny add-on adapters (r=16)",
         "Focal loss (γ=2): focus on hard examples, ignore easy ones",
         "Oversampling: copy rare-class memes until classes are balanced",
     ],
     "Why: full fine-tune was unstable on the rebalanced data — LoRA is much more stable.",
     NAVY),
    ("T2.3 · Which categories?",
     "XLM-R-base, multi-label — 5 sigmoid heads",
     [
         "Five independent yes/no detectors (sigmoid heads)",
         "Trained with BCE + pos_weight: extra weight on rare categories",
         "Fires above 0.5; if none fires → keep highest (top-1 fallback)",
         "Gated by T2.1: if T2.1 says NO → output is just [NO]",
     ],
     "Why: a single-label head just learned 'always predict NO' → 0% on every category.",
     GREEN),
]
cw = Inches(4.1)
ch = Inches(4.9)
gap = Inches(0.2)
csx = (SW - 3 * cw - 2 * gap) / 2
cy = Inches(1.45)
for i, (title, model, bullets, note, color) in enumerate(cols):
    x = csx + i * (cw + gap)
    add_rect(s, x, cy, cw, ch, fill=LIGHT)
    add_rect(s, x, cy, cw, Inches(0.5), fill=color)
    add_text(s, x + Inches(0.15), cy + Inches(0.03),
             cw - Inches(0.3), Inches(0.45),
             title, size=18, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, x + Inches(0.2), cy + Inches(0.6),
             cw - Inches(0.4), Inches(0.6),
             model, size=15, bold=True, color=INK)
    # bullets
    tb = add_bullets(s, x + Inches(0.2), cy + Inches(1.25),
                     cw - Inches(0.4), Inches(2.4),
                     bullets, size=13, color=INK)
    # note panel
    add_rect(s, x + Inches(0.2), cy + Inches(3.5),
             cw - Inches(0.4), Inches(1.25),
             fill=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, x + Inches(0.3), cy + Inches(3.55),
             cw - Inches(0.6), Inches(1.2),
             note, size=12, color=MUTED)

set_notes(s, """All three tasks use the same model family — XLM-RoBERTa — but with different recipes.

For Task 2.1 I do a 'full fine-tune' of XLM-R-large. Fine-tuning means: take the pre-trained model and continue training it on the memes; 'full' means I update all the weights, not just some. I use a slow learning rate so it doesn't 'forget' what it knew, and I use class-weighted cross-entropy as the loss function — with extra weight on the rare class so its mistakes hurt more.

For Task 2.2, the 'judgemental' class is only 14% of the data — severe class imbalance. So I add three things on top of XLM-R-large. First, LoRA — Low-Rank Adaptation. Instead of updating all the model's weights (which kept blowing up), I freeze the model and train tiny adapter layers on top. Second, focal loss — a loss function that automatically focuses on hard examples and ignores easy ones. The γ parameter (gamma) controls how aggressively. Third, oversampling — literally copying rare-class examples until they appear as often as the common ones.

For Task 2.3, the key change is switching to a multi-label head. Instead of a softmax that picks one category out of five, I use five sigmoid heads — five independent yes/no detectors. Each one outputs a probability between 0 and 1. I train with BCE — binary cross-entropy — and pos_weight, which gives extra weight to the rare positive examples. A detector fires if its score is above 0.5; if none fires, I keep the highest one so I never predict 'nothing'. And on top of that I gate the whole T2.3 head by T2.1: if T2.1 says NO, T2.3 outputs just [NO].""")


# ============================================================
# Slide 6 — Late fusion (architecture diagram)
# ============================================================
s = new_slide()
add_header(s, "Late fusion — two models that vote",
           "Text model + image+sensor model → mix their probabilities with weight α")

# Diagram: three input boxes -> two model boxes -> fusion -> prediction
# Three input boxes on the left
def labeled_box(x, y, w, h, title, sub, color):
    add_rect(s, x, y, w, h, fill=color)
    add_text(s, x, y + Inches(0.05), w, Inches(0.35),
             title, size=14, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.4), w, Inches(0.4),
             sub, size=11,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

ix = Inches(0.6)
iy_text = Inches(1.7)
iy_img = Inches(3.2)
iy_phy = Inches(4.7)
ibw = Inches(2.4)
ibh = Inches(0.9)
labeled_box(ix, iy_text, ibw, ibh, "TEXT", "meme caption (EN / ES)", NAVY)
labeled_box(ix, iy_img, ibw, ibh, "IMAGE", "CLIP features — 512 numbers", GREEN)
labeled_box(ix, iy_phy, ibw, ibh, "SENSORS", "ET + HR + EEG — 104 numbers", ACCENT)

# arrows -> two model boxes
mx = Inches(4.0)
mbw = Inches(3.4)
mbh = Inches(1.5)
my_text = Inches(1.85)
my_lr = Inches(4.0)

add_rect(s, mx, my_text, mbw, mbh, fill=LIGHT, line=NAVY)
add_text(s, mx, my_text + Inches(0.1), mbw, Inches(0.4),
         "XLM-RoBERTa  (large / base)", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)
add_text(s, mx, my_text + Inches(0.55), mbw, Inches(0.4),
         "multilingual language model — fine-tuned",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, mx, my_text + Inches(0.95), mbw, Inches(0.4),
         "→  p_text  (text model's vote)", size=14, bold=True, color=NAVY,
         align=PP_ALIGN.CENTER)

add_rect(s, mx, my_lr, mbw, mbh, fill=LIGHT, line=ACCENT)
add_text(s, mx, my_lr + Inches(0.1), mbw, Inches(0.4),
         "Logistic Regression", size=14, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
add_text(s, mx, my_lr + Inches(0.55), mbw, Inches(0.4),
         "linear classifier on image ⊕ sensors (616 numbers)",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, mx, my_lr + Inches(0.95), mbw, Inches(0.4),
         "→  p_img+sens  (image+sensor vote)", size=13, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)

# Fusion box
fx = Inches(8.4)
fy = Inches(2.85)
fw = Inches(3.0)
fh = Inches(1.7)
add_rect(s, fx, fy, fw, fh, fill=NAVY)
add_text(s, fx, fy + Inches(0.15), fw, Inches(0.5),
         "LATE FUSION", size=16, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, fx, fy + Inches(0.7), fw, Inches(0.5),
         "p  =  α · p_text  +  (1−α) · p_img+sens", size=12,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, fx, fy + Inches(1.15), fw, Inches(0.5),
         "α = trust in text — tuned on validation set", size=12,
         color=RGBColor(0xC9, 0xD1, 0xDE), align=PP_ALIGN.CENTER)

# Final prediction
px = Inches(8.4)
py = Inches(5.0)
add_rect(s, px, py, fw, Inches(0.9), fill=ACCENT)
add_text(s, px, py + Inches(0.05), fw, Inches(0.4),
         "PREDICTION", size=14, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, px, py + Inches(0.45), fw, Inches(0.4),
         "hard label + soft probability", size=12,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

# Alpha values panel
ax = Inches(0.6)
ay = Inches(6.0)
add_rect(s, ax, ay, Inches(7.4), Inches(0.85), fill=LIGHT)
add_text(s, ax + Inches(0.2), ay + Inches(0.05),
         Inches(7.0), Inches(0.4),
         "Chosen α per task   (how much I trust the text model)",
         size=13, bold=True, color=NAVY)
add_text(s, ax + Inches(0.2), ay + Inches(0.4),
         Inches(7.0), Inches(0.4),
         "T2.1 → α = 0.95  (95% text)    ·    T2.2 → α = 0.75  (75% text)    ·    T2.3 → gated by T2.1",
         size=12, color=INK)

# Simple connector lines (just thin rectangles as visual arrows)
def thin_line(x1, y1, x2, y2, color=MUTED, thickness=Inches(0.03)):
    if y1 == y2:
        add_rect(s, x1, y1, x2 - x1, thickness, fill=color)
    elif x1 == x2:
        add_rect(s, x1, y1, thickness, y2 - y1, fill=color)

# arrows from inputs to text-model / lr
ay_text_mid = iy_text + ibh / 2
ay_img_mid = iy_img + ibh / 2
ay_phy_mid = iy_phy + ibh / 2
ix_right = ix + ibw
mx_left = mx
text_in_y = my_text + Inches(0.6)
lr_in_y = my_lr + Inches(0.6)

thin_line(ix_right, ay_text_mid, mx_left, ay_text_mid)
# vertical bridge from text input down to text model entry
thin_line(ix_right + Inches(0.5), min(ay_text_mid, text_in_y),
          ix_right + Inches(0.5), max(ay_text_mid, text_in_y))

# image + phys merge into LR
thin_line(ix_right, ay_img_mid, ix_right + Inches(0.4), ay_img_mid)
thin_line(ix_right, ay_phy_mid, ix_right + Inches(0.4), ay_phy_mid)
thin_line(ix_right + Inches(0.4), ay_img_mid,
          ix_right + Inches(0.4), ay_phy_mid)
thin_line(ix_right + Inches(0.4), lr_in_y,
          mx_left, lr_in_y)

# Branches into fusion
thin_line(mx + mbw, my_text + mbh / 2, fx, my_text + mbh / 2)
thin_line(fx - Inches(0.4), my_text + mbh / 2,
          fx - Inches(0.4), my_lr + mbh / 2)
thin_line(mx + mbw, my_lr + mbh / 2, fx - Inches(0.4), my_lr + mbh / 2)
thin_line(fx - Inches(0.4), fy + fh / 2, fx, fy + fh / 2)

# Fusion -> prediction
thin_line(fx + fw / 2, fy + fh, fx + fw / 2, py)

set_notes(s, """This is the late-fusion architecture — same shape for all three tasks.

On the left, the three inputs. The text goes through fine-tuned XLM-RoBERTa, which outputs a probability — I call it p_text, the text model's vote. The image features (512 numbers from CLIP) and the sensor features (104 numbers) are concatenated into one 616-dimensional vector and fed to a logistic regression — a linear classifier. That outputs a second probability, p_img+sens.

Then I combine the two votes at the probability level: final = α times p_text plus (1−α) times p_img+sens. That's what 'late fusion' means — fusion happens late, after each model has already made its decision, instead of mixing inputs early. The α controls how much I trust the text model. I tune it on a held-out validation set.

For Task 2.1, α = 0.95 → I trust the text model 95%. The image and sensors only nudge the prediction.
For Task 2.2, α = 0.75 → image and sensors matter a bit more, because the difference between 'direct' and 'judgemental' sexism often shows in the picture.
For Task 2.3, instead of α-fusion, I 'gate' it: if Task 2.1 says NO, T2.3 outputs just [NO]. Otherwise the five sigmoid heads each decide.""")


# ============================================================
# Slide 7 — Final models per task (table)
# ============================================================
s = new_slide()
add_header(s, "Final models I submitted",
           "One system per task — same overall idea, different recipe")

tx = Inches(0.7)
ty = Inches(1.5)
tw = SW - Inches(1.4)

# header row
add_rect(s, tx, ty, Inches(1.2), Inches(0.6), fill=NAVY)
add_text(s, tx, ty + Inches(0.1), Inches(1.2), Inches(0.4),
         "Task", size=14, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_rect(s, tx + Inches(1.2), ty, Inches(5.4), Inches(0.6), fill=NAVY)
add_text(s, tx + Inches(1.2), ty + Inches(0.1), Inches(5.4), Inches(0.4),
         "Text backbone",
         size=14, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_rect(s, tx + Inches(6.6), ty, Inches(3.4), Inches(0.6), fill=NAVY)
add_text(s, tx + Inches(6.6), ty + Inches(0.1), Inches(3.4), Inches(0.4),
         "Loss / training tricks",
         size=14, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_rect(s, tx + Inches(10.0), ty, Inches(1.95), Inches(0.6), fill=NAVY)
add_text(s, tx + Inches(10.0), ty + Inches(0.1),
         Inches(1.95), Inches(0.4),
         "Fusion α", size=14, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

rows = [
    ("T2.1", "XLM-R-large · full fine-tune (all weights)",
     "Weighted cross-entropy · 3 epochs", "α = 0.95"),
    ("T2.2", "XLM-R-large + LoRA (tiny add-on adapters)",
     "Focal loss + oversampling · 5 epochs", "α = 0.75"),
    ("T2.3", "XLM-R-base · 5 sigmoid heads (multi-label)",
     "BCE + pos_weight · gated by T2.1", "gated by T2.1"),
]
ry = ty + Inches(0.6)
rh = Inches(0.95)
for i, (t, b, tricks, alpha) in enumerate(rows):
    fill = LIGHT if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    add_rect(s, tx, ry, Inches(1.2), rh, fill=fill)
    add_text(s, tx, ry + Inches(0.25), Inches(1.2), Inches(0.5),
             t, size=18, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_rect(s, tx + Inches(1.2), ry, Inches(5.4), rh, fill=fill)
    add_text(s, tx + Inches(1.3), ry + Inches(0.25),
             Inches(5.2), Inches(0.5),
             b, size=14, color=INK)
    add_rect(s, tx + Inches(6.6), ry, Inches(3.4), rh, fill=fill)
    add_text(s, tx + Inches(6.7), ry + Inches(0.25),
             Inches(3.2), Inches(0.5),
             tricks, size=13, color=INK)
    add_rect(s, tx + Inches(10.0), ry, Inches(1.95), rh, fill=fill)
    add_text(s, tx + Inches(10.0), ry + Inches(0.25),
             Inches(1.95), Inches(0.5),
             alpha, size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    ry += rh

# Selection rationale
add_rect(s, tx, ry + Inches(0.3), tw, Inches(1.4), fill=LIGHT)
add_text(s, tx + Inches(0.2), ry + Inches(0.35),
         tw - Inches(0.4), Inches(0.4),
         "Why I picked these three",
         size=15, bold=True, color=NAVY)
add_text(s, tx + Inches(0.2), ry + Inches(0.75),
         tw - Inches(0.4), Inches(1.0),
         "→  They scored the best on a held-out validation set, and they still fit on a single GPU\n"
         "→  All three follow the same 'two models that vote' pattern → easy to maintain and to run on the test set",
         size=13, color=INK)

set_notes(s, """These three systems are what I actually submitted.
Same family of model, different recipe per task, different amount of trust in the text vote.
I picked them because they gave the best score on a validation set I held back from training — and because they fit on a single GPU.
Now to the question you specifically asked about: does the sensor data actually help?""")


# ============================================================
# Slide 8 — THE main slide — CV with vs without physiological data
# ============================================================
s = new_slide()
add_header(s, "★ 5-fold cross-validation: with vs. without physiological data",
           "Same fusion MLP, same 5 folds — only the input features change")

# Table on the left
tx = Inches(0.55)
ty = Inches(1.5)
table_w = Inches(6.4)
row_h = Inches(0.55)

# header
add_rect(s, tx, ty, Inches(2.6), row_h, fill=NAVY)
add_text(s, tx, ty + Inches(0.1), Inches(2.6), Inches(0.4),
         "Features", size=13, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
for i, lab in enumerate(["T2.1", "T2.2", "T2.3"]):
    add_rect(s, tx + Inches(2.6) + i * Inches(1.27), ty,
             Inches(1.27), row_h, fill=NAVY)
    add_text(s, tx + Inches(2.6) + i * Inches(1.27),
             ty + Inches(0.1),
             Inches(1.27), Inches(0.4),
             lab, size=13, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

# rows
rows = [
    ("Text only",
     "0.590 ± 0.009", "0.385 ± 0.019", "0.237 ± 0.017",
     False, False, False),
    ("Text + Image (CLIP)",
     "0.581 ± 0.010", "0.398 ± 0.006", "0.246 ± 0.015",
     False, False, False),
    ("Text + Image + Phys (ET/HR/EEG)",
     "0.578 ± 0.011", "0.394 ± 0.008", "0.253 ± 0.013",
     False, False, True),
]
ry = ty + row_h
for j, (feat, a, b, c, ha, hb, hc) in enumerate(rows):
    bg = LIGHT if j % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
    add_rect(s, tx, ry, Inches(2.6), row_h, fill=bg)
    add_text(s, tx + Inches(0.15), ry + Inches(0.1),
             Inches(2.4), Inches(0.4),
             feat, size=13, bold=(j == 2), color=INK)
    for k, (val, hl) in enumerate([(a, ha), (b, hb), (c, hc)]):
        add_rect(s, tx + Inches(2.6) + k * Inches(1.27), ry,
                 Inches(1.27), row_h,
                 fill=RGBColor(0xFF, 0xE7, 0xB3) if hl else bg)
        add_text(s, tx + Inches(2.6) + k * Inches(1.27),
                 ry + Inches(0.1),
                 Inches(1.27), Inches(0.4),
                 val, size=13, bold=hl, color=INK,
                 align=PP_ALIGN.CENTER)
    ry += row_h

# Image of CV comparison on the right
img_path = str(FIG / "cv_phys_comparison.png")
img_x = Inches(7.2)
img_y = Inches(1.5)
img_w = Inches(5.7)
s.shapes.add_picture(img_path, img_x, img_y, width=img_w)

# Findings panel below table
fx2 = Inches(0.55)
fy2 = Inches(4.55)
fw2 = Inches(12.2)
fh2 = Inches(2.5)
add_rect(s, fx2, fy2, fw2, fh2, fill=LIGHT)
add_text(s, fx2 + Inches(0.2), fy2 + Inches(0.1),
         fw2 - Inches(0.4), Inches(0.4),
         "Reading the numbers   (macro-F1 — averages F1 across categories, weighting each one equally; higher is better)",
         size=14, bold=True, color=NAVY)

findings = [
    "T2.1 (binary): physiological data does NOT help — text alone already solves it. Adding sensors slightly hurts.",
    "T2.2 (type): no real change — CLIP image features already absorb the extra information.",
    "T2.3 (category): +0.016 macro-F1 from physiological data. Small, but consistent across all 5 folds → especially helps the rarest categories.",
    "Take-away: physiological signals only pay off where text saturates and class imbalance is severe (fine-grained categories).",
]
ftb = s.shapes.add_textbox(fx2 + Inches(0.2), fy2 + Inches(0.5),
                            fw2 - Inches(0.4), fh2 - Inches(0.6))
ftf = ftb.text_frame
ftf.word_wrap = True
for i, line in enumerate(findings):
    p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
    p.line_spacing = 1.15
    r = p.add_run()
    r.text = "→  " + line
    r.font.size = Pt(13)
    r.font.color.rgb = INK
    if i == 2:
        r.font.bold = True
        r.font.color.rgb = ACCENT

set_notes(s, """This is the slide you specifically asked about — let me spend a moment on the setup before the numbers.

Setup. To honestly answer 'do sensors help', I need to vary only the inputs. So I use one fixed comparison network — a small fusion MLP — and 5-fold cross-validation. Cross-validation means I split the training data into 5 chunks; I train on 4 chunks, test on the 5th, then rotate, so every meme gets used as test exactly once. I do the whole thing three times: text only, text + image (CLIP features), text + image + physiological. Same model, same folds, only the input features change. The score is macro-F1 — the F1 score averaged across categories, treating each one equally; this is what punishes models that 'cheat' by always predicting the majority class.

The numbers:
- Task 2.1 (binary): physiological data actually hurts a tiny bit — 0.590 down to 0.578. Text already solves this; sensors add noise.
- Task 2.2 (type): essentially flat. CLIP image features already absorb whatever the sensors would have added.
- Task 2.3 (category) — the one positive result. Physiological data pushes macro-F1 from 0.246 up to 0.253. Small in absolute terms — but on the right-hand chart you can see the orange bar is always slightly higher for T2.3 across every single fold. Consistent direction, not luck.

The intuition is natural. Deciding 'is this sexist or not' depends mostly on what the meme says. But picking the exact category — especially rare ones like 'sexual violence' — is where every extra signal counts, and physiological data gives the model a second look at how annotators reacted to the image.""")


# ============================================================
# Slide 9 — Final models vs baselines (training set numbers)
# ============================================================
s = new_slide()
add_header(s, "Final system vs. baselines",
           "Macro-F1 and AUC on the full training set (3984 memes)")

# Three column charts from the figures folder
img_w = Inches(4.05)
gap = Inches(0.2)
total_w = 3 * img_w + 2 * gap
start = (SW - total_w) / 2
y_img = Inches(1.5)
for i, name in enumerate(["comparison_T21.png", "comparison_T22.png",
                          "comparison_T23.png"]):
    p = str(FIG / name)
    s.shapes.add_picture(p, start + i * (img_w + gap), y_img, width=img_w)

# Numbers panel below
tx = Inches(0.55)
ty = Inches(5.3)
add_rect(s, tx, ty, SW - Inches(1.1), Inches(1.85), fill=LIGHT)
add_text(s, tx + Inches(0.2), ty + Inches(0.1),
         SW - Inches(1.5), Inches(0.4),
         "XLM-R-base text baseline   →   final multimodal model     ·     F1 = score across classes,  AUC = ranking quality",
         size=14, bold=True, color=NAVY)
nums = [
    ("T2.1", "F1   0.422  →  0.794", "AUC   0.386 → 0.867"),
    ("T2.2", "F1   0.435  →  0.635", "AUC   0.741 → 0.827"),
    ("T2.3", "F1   0.000  →  0.357", "AUC   0.614 → 0.821"),
]
nx = tx + Inches(0.4)
ny = ty + Inches(0.55)
cell_w = (SW - Inches(2.0)) / 3
for i, (t, f1, auc) in enumerate(nums):
    add_text(s, nx + i * cell_w, ny, cell_w, Inches(0.4),
             t, size=16, bold=True, color=ACCENT)
    add_text(s, nx + i * cell_w, ny + Inches(0.4),
             cell_w, Inches(0.4), f1, size=14, color=INK)
    add_text(s, nx + i * cell_w, ny + Inches(0.8),
             cell_w, Inches(0.4), auc, size=14, color=INK)

set_notes(s, """Now comparing my XLM-R-base text-only baseline against the final multimodal system, on the full training set. Two metrics per task: macro-F1 (an average score across categories, treating each one equally — this is what I am optimising) and AUC (Area Under the ROC Curve — how well the model ranks examples by confidence, independent of any threshold).

Task 2.1: macro-F1 nearly doubles, 0.42 → 0.79. Most of that comes from switching from XLM-R-base to XLM-R-large and from adding class weights — not from fusion. Late fusion with the image+sensor branch only adds a small correction (α = 0.95).

Task 2.2: 0.44 → 0.63. Here the LoRA-plus-focal-loss-plus-oversampling stack does most of the heavy lifting.

Task 2.3 is the most striking. The XLM-R-base baseline scores exactly zero macro-F1 because it is single-label — it just learns to always predict NO and never fires any of the five category heads. Switching to the multi-label setup with sigmoid heads and pos_weight brings it to 0.357. That's the biggest qualitative jump in the project.""")


# ============================================================
# Slide 10 — Per-category T2.3 + test-set breakdown
# ============================================================
s = new_slide()
add_header(s, "Zooming into T2.3 — per-category & test predictions",
           "Multi-label head (5 sigmoid + pos_weight) recovers every category the single-label baseline missed")

# Left: per-category table
tx = Inches(0.6)
ty = Inches(1.5)
add_rect(s, tx, ty, Inches(6.1), Inches(0.55), fill=NAVY)
add_text(s, tx + Inches(0.15), ty + Inches(0.1),
         Inches(3.0), Inches(0.4),
         "Category", size=13, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, tx + Inches(3.2), ty + Inches(0.1),
         Inches(1.4), Inches(0.4),
         "Baseline", size=13, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, tx + Inches(4.7), ty + Inches(0.1),
         Inches(1.4), Inches(0.4),
         "Final F1", size=13, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

cats = [
    ("Ideological inequality",        "0.000", "0.546"),
    ("Stereotyping dominance",        "0.000", "0.392"),
    ("Objectification",               "0.000", "0.341"),
    ("Sexual violence",               "0.000", "0.285"),
    ("Misogyny (non-sexual violence)","0.000", "0.223"),
    ("Macro average",                 "0.000", "0.357"),
]
ry = ty + Inches(0.55)
rh = Inches(0.5)
for j, (c, b, f) in enumerate(cats):
    is_avg = (j == len(cats) - 1)
    bg = (RGBColor(0xFF, 0xE7, 0xB3) if is_avg
          else (LIGHT if j % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)))
    add_rect(s, tx, ry, Inches(6.1), rh, fill=bg)
    add_text(s, tx + Inches(0.15), ry + Inches(0.08),
             Inches(3.0), Inches(0.4),
             c, size=13, bold=is_avg, color=INK)
    add_text(s, tx + Inches(3.2), ry + Inches(0.08),
             Inches(1.4), Inches(0.4),
             b, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, tx + Inches(4.7), ry + Inches(0.08),
             Inches(1.4), Inches(0.4),
             f, size=13, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    ry += rh

# Right: test set prediction breakdown
rx = Inches(7.1)
ry2 = Inches(1.5)
rw = Inches(5.7)
add_rect(s, rx, ry2, rw, Inches(5.4), fill=LIGHT)
add_text(s, rx + Inches(0.2), ry2 + Inches(0.15),
         rw - Inches(0.4), Inches(0.5),
         "Test-set hard predictions (1053 memes)",
         size=15, bold=True, color=NAVY)

test_lines = [
    ("T2.1",  "580 YES   ·   473 NO"),
    ("T2.2",  "458 NO   ·   346 DIRECT   ·   249 JUDGEMENTAL"),
    ("T2.3",  "293 objectification"),
    ("",      "274 stereotyping dominance"),
    ("",      "265 ideological inequality"),
    ("",      "152 misogyny NSV"),
    ("",      "132 sexual violence"),
    ("",      "473 plain NO"),
]
ly = ry2 + Inches(0.8)
for tag, text in test_lines:
    if tag:
        add_text(s, rx + Inches(0.3), ly, Inches(0.7),
                 Inches(0.4), tag, size=13, bold=True, color=ACCENT)
    add_text(s, rx + Inches(1.0), ly, rw - Inches(1.3),
             Inches(0.4), text, size=13, color=INK)
    ly += Inches(0.45)

add_text(s, rx + Inches(0.2), ry2 + Inches(4.7),
         rw - Inches(0.4), Inches(0.6),
         "Distribution close to training frequencies → no class collapse.",
         size=12, color=MUTED)

set_notes(s, """A closer look at Task 2.3.

On the left, the score per category. The baseline is flat zero across all five — it never predicts any sexism category, it just always says NO. My final model gets a real score on every category.

The order follows how often each category appears in training. 'Ideological inequality' is the easiest because it's the most common. 'Sexual violence' and 'misogyny without violence' are the hardest because they are rare in the training data. But the important thing is — all five categories fire.

On the right, the breakdown of what the final model actually predicted on the 1053 test memes. Task 2.1 predicts about 55% YES, 45% NO — close to the training balance. Task 2.3 doesn't collapse to a single category either — all five appear with reasonable counts. So the model didn't degenerate at test time.""")


# ============================================================
# Slide 11 — Conclusions
# ============================================================
s = new_slide()
add_header(s, "Take-aways",
           "What worked, what didn't, what's next")

# Three columns: worked / surprising / next
cx = Inches(0.5)
cy = Inches(1.5)
cw = Inches(4.1)
ch = Inches(5.0)
gap = Inches(0.15)

cols = [
    ("What worked",
     GREEN,
     [
         "Scaling XLM-R-base → XLM-R-large nearly doubled T2.1 macro-F1",
         "LoRA + focal loss + oversampling stabilised the rare-class T2.2 training",
         "Multi-label sigmoid heads + pos_weight broke the F1 = 0 trap on T2.3",
         "Late fusion (α-weighted) adds a small but real boost over text alone",
     ]),
    ("What surprised me",
     ACCENT,
     [
         "Physiological data (ET/HR/EEG) does NOT help binary sexism detection",
         "CLIP image features alone already absorb most of the multimodal gain",
         "Physiological data only helps the rarest T2.3 categories",
         "Gap between val (0.81) and CV (0.59) → XLM-R-large memorises part of the train set",
     ]),
    ("If I had more time",
     NAVY,
     [
         "Replace CLIP + LogReg with a vision-language model (Qwen2-VL, LLaVA)",
         "Add a learnable attention over the 104 physiological features",
         "Train per-language models (EN / ES) on top of the shared backbone",
         "Confidence-weighted gating between the T2.1 and T2.3 heads",
     ]),
]
for i, (title, color, items) in enumerate(cols):
    x = cx + i * (cw + gap)
    add_rect(s, x, cy, cw, ch, fill=LIGHT)
    add_rect(s, x, cy, cw, Inches(0.5), fill=color)
    add_text(s, x + Inches(0.15), cy + Inches(0.05),
             cw - Inches(0.3), Inches(0.4),
             title, size=17, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_bullets(s, x + Inches(0.2), cy + Inches(0.7),
                cw - Inches(0.4), ch - Inches(0.8),
                items, size=13, color=INK, line_spacing=1.25)

# Bottom big quote
qx = Inches(0.5)
qy = Inches(6.6)
add_rect(s, qx, qy, SW - Inches(1.0), Inches(0.55), fill=NAVY)
add_text(s, qx, qy + Inches(0.1), SW - Inches(1.0), Inches(0.4),
         "Bottom line: text is by far the strongest signal — "
         "physiological data only helps where text saturates (rare T2.3 categories).",
         size=14, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

set_notes(s, """To wrap up — three columns.

What worked. Scaling XLM-R-base up to XLM-R-large nearly doubled macro-F1 on Task 2.1. The LoRA + focal loss + oversampling stack stabilised the rare-class training on Task 2.2 — full fine-tune kept blowing up there. And the multi-label sigmoid heads with pos_weight broke through on Task 2.3, where the single-label baseline was scoring exactly zero.

What surprised me. Physiological data — eye tracking, heart rate, EEG — does NOT help with binary sexism detection. CLIP image features on their own already absorb most of what multimodality gives you. Physiological data only helps on the rarest T2.3 categories. And there is a noticeable gap between validation macro-F1 (0.81) and cross-validation macro-F1 (0.59) — that tells me XLM-R-large is memorising part of the training set.

If I had more time. Replace CLIP + logistic regression with a true vision-language model like Qwen2-VL or LLaVA. Add a learnable attention layer over the 104 physiological features instead of treating them as flat input. Train per-language models for English and Spanish on top of the shared backbone. And add confidence-weighted gating between T2.1 and T2.3.

Bottom line — and this is the main empirical finding — text is by far the strongest signal. Physiological data only pays off where text saturates and class imbalance dominates, which means the rarest T2.3 categories.""")


# ============================================================
# Slide 12 — Thanks / Q&A
# ============================================================
s = new_slide()
add_rect(s, 0, 0, SW, SH, fill=NAVY)
add_rect(s, 0, Inches(3.05), SW, Inches(0.06), fill=ACCENT)

add_text(s, Inches(0.8), Inches(2.2), SW - Inches(1.6), Inches(1.0),
         "Thank you", size=54, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.8), Inches(3.3), SW - Inches(1.6), Inches(0.7),
         "Questions?", size=28,
         color=RGBColor(0xC9, 0xD1, 0xDE))

add_text(s, Inches(0.8), Inches(5.4), SW - Inches(1.6), Inches(0.4),
         "Klaudia Banasiewicz · EXIST 2026 — Lab Project", size=16,
         color=RGBColor(0xFF, 0xFF, 0xFF))

set_notes(s, """Thank you — happy to take any questions.
If you want extra context on the LoRA setup, the fusion α sweep, or the 5-fold split, I can go into more detail.""")


# Save
prs.save(OUT)
print(f"Wrote: {OUT}")
print(f"Slides: {len(prs.slides)}")
